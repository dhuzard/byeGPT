"""
backend/app/main.py — FastAPI entry point and API routes for byeGPT Studio.
"""

from __future__ import annotations

import asyncio
import base64
import binascii
import json
import logging
import os
import re
import tempfile
import uuid
import wave
from pathlib import Path
from typing import Any
from io import BytesIO

from fastapi import BackgroundTasks, FastAPI, File, Form, HTTPException, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
from pydantic import BaseModel, Field

from .auth_manager import is_authenticated, login_and_save
from .cloud import (
    NotebookUploadError,
    batch_upload,
    generate_audio_overview,
    generate_mind_map,
    generate_quiz,
    generate_slides,
    revise_slide,
)
from .jobs import jobs
from .parser import convert_with_anchors, inject_context_anchors
from .storage import (
    StorageManager,
    build_csv,
    build_markdown_table,
    derive_slide_table_rows,
)
from .topics import build_topic_laboratory
from byegpt.formatter import format_conversation
from byegpt.parser import load_conversations
from byegpt.taxonomy import build_taxonomy, conversation_uid

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app = FastAPI(
    title="byeGPT Studio API",
    description=(
        "Backend for the byeGPT Studio dashboard. Converts ChatGPT exports, "
        "manages NotebookLM notebooks, and persists generated artifacts."
    ),
    version="4.0.0",
)

_ALLOWED_ORIGINS = os.environ.get(
    "CORS_ORIGINS",
    "http://localhost:5173,http://localhost:3000",
).split(",")

app.add_middleware(
    CORSMiddleware,
    allow_origins=_ALLOWED_ORIGINS,
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

_STORAGE_DIR = Path(os.environ.get("BYEGPT_STORAGE", ".byegpt"))
_COOKIES_PATH = _STORAGE_DIR / "storage.json"
_CONVERTED_DIR = _STORAGE_DIR / "converted"
_TUSD_UPLOAD_DIR = Path(os.environ.get("BYEGPT_TUSD_UPLOAD_DIR", ".uploads"))
_SEARCH_DIR = _STORAGE_DIR / "index"
_DEMO_MODE = os.environ.get("BYEGPT_DEMO_MODE", "false").lower() in {"1", "true", "yes"}
_storage = StorageManager(_STORAGE_DIR)


def _get_storage() -> StorageManager:
    return _storage


class UploadNotebookRequest(BaseModel):
    notebook_title: str = "byeGPT Archive"
    output_dir: str = str(_CONVERTED_DIR)
    passport_id: str | None = None


class ArtifactJobRequest(BaseModel):
    types: list[str] = Field(default_factory=lambda: ["mind_map", "audio", "slides", "quiz"])


class ExportArtifactRequest(BaseModel):
    format: str


class ReviseSlideRequest(BaseModel):
    revision_prompt: str
    artifact_id: str


class SearchIndexRequest(BaseModel):
    input_dir: str


class SearchQueryRequest(BaseModel):
    text: str
    n_results: int = 5


class NotebookSelection(BaseModel):
    category: str
    subcategory: str


class DerivedNotebookRequest(BaseModel):
    title: str
    passport_id: str
    parent_output_dir: str | None = None
    parent_notebook_id: str | None = None
    selections: list[NotebookSelection] = Field(default_factory=list)


class TusUploadInfo(BaseModel):
    ID: str
    Size: int
    Offset: int
    MetaData: dict[str, str] = Field(default_factory=dict)


@app.get("/health", tags=["system"])
async def health() -> dict[str, str]:
    return {"status": "ok", "version": "4.0.0"}


@app.get("/auth/status", tags=["auth"])
async def auth_status() -> dict[str, bool]:
    return {"authenticated": _DEMO_MODE or is_authenticated(_COOKIES_PATH)}


@app.post("/auth/login", tags=["auth"])
async def auth_login(background_tasks: BackgroundTasks) -> dict[str, str]:
    if _DEMO_MODE:
        return {"status": "demo_mode", "message": "Demo mode is enabled. NotebookLM login is bypassed."}

    if is_authenticated(_COOKIES_PATH):
        return {"status": "already_authenticated"}

    if os.name != "nt" and not os.environ.get("DISPLAY"):
        raise HTTPException(
            status_code=501,
            detail=(
                "Interactive NotebookLM login is unavailable in Docker without an X server. "
                "Provide a valid .byegpt/storage.json session file or run the backend outside Docker for login."
            ),
        )

    background_tasks.add_task(login_and_save, _COOKIES_PATH, headless=False)
    return {"status": "login_started", "message": "Open the browser window to complete login."}


@app.post("/convert", tags=["convert"])
async def convert_export(
    file: UploadFile = File(..., description="ChatGPT .zip or conversations.json export"),
    max_size_mb: float = Form(7.0),
    include_thinking: bool = Form(True),
    include_attachments: bool = Form(True),
) -> dict[str, Any]:
    suffix = Path(file.filename or "upload").suffix or ".zip"
    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
        tmp.write(await file.read())
        tmp_path = Path(tmp.name)

    output_dir = _CONVERTED_DIR / str(uuid.uuid4())
    try:
        return await _convert_archive(
            input_path=tmp_path,
            output_dir=output_dir,
            max_size_mb=max_size_mb,
            include_thinking=include_thinking,
            include_attachments=include_attachments,
        )
    finally:
        tmp_path.unlink(missing_ok=True)


@app.post("/convert/tus/{upload_id}", tags=["convert"])
async def convert_tus_upload(
    upload_id: str,
    max_size_mb: float = 7.0,
    include_thinking: bool = True,
    include_attachments: bool = True,
) -> dict[str, Any]:
    info = _load_tus_upload_info(upload_id)
    source_path = _TUSD_UPLOAD_DIR / upload_id
    if not source_path.exists():
        raise HTTPException(status_code=404, detail=f"Tus upload '{upload_id}' not found.")

    if source_path.stat().st_size < info.Size:
        raise HTTPException(status_code=409, detail="Upload is not complete yet.")

    metadata = {key: _normalize_tusd_metadata(value) for key, value in info.MetaData.items()}
    original_name = metadata.get("filename")
    if not original_name:
        raise HTTPException(status_code=400, detail="Tus upload metadata is missing filename.")

    suffix = Path(original_name).suffix.lower()
    if suffix not in {".zip", ".json"}:
        raise HTTPException(status_code=400, detail="Unsupported upload type. Use .zip or .json.")

    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
        tmp.write(source_path.read_bytes())
        tmp_path = Path(tmp.name)

    output_dir = _CONVERTED_DIR / str(uuid.uuid4())
    try:
        return await _convert_archive(
            input_path=tmp_path,
            output_dir=output_dir,
            max_size_mb=max_size_mb,
            include_thinking=include_thinking,
            include_attachments=include_attachments,
        )
    finally:
        tmp_path.unlink(missing_ok=True)


@app.post("/persona", tags=["convert"])
async def generate_passport(
    file: UploadFile = File(..., description="ChatGPT .zip or conversations.json export"),
) -> dict[str, Any]:
    from core.persona import build_passport_bundle

    suffix = Path(file.filename or "upload").suffix or ".zip"
    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
        tmp.write(await file.read())
        tmp_path = Path(tmp.name)

    try:
        bundle = await asyncio.to_thread(build_passport_bundle, tmp_path)
    finally:
        tmp_path.unlink(missing_ok=True)

    passport = _get_storage().save_passport_bundle(
        markdown=bundle["markdown"],
        taxonomy=bundle["taxonomy"],
    )
    return {
        "passport_markdown": bundle["markdown"],
        "passport_id": passport["passport_id"],
        "saved_path": passport["path"],
        "taxonomy": bundle["taxonomy"],
    }


@app.post("/notebooks/upload", tags=["notebooklm"])
async def upload_to_notebooklm(body: UploadNotebookRequest) -> dict[str, Any]:
    if not _DEMO_MODE and not is_authenticated(_COOKIES_PATH):
        raise HTTPException(status_code=401, detail="Not authenticated. Call /auth/login first.")

    output_dir = Path(body.output_dir)
    if not output_dir.exists():
        raise HTTPException(status_code=404, detail=f"Output dir '{output_dir}' not found.")

    md_files = sorted(output_dir.rglob("*.md"))
    if not md_files:
        raise HTTPException(status_code=404, detail="No Markdown files found in output_dir.")

    if _DEMO_MODE:
        chunks = [md_files[index : index + 50] for index in range(0, len(md_files), 50)]
        notebook_ids = [f"demo_{uuid.uuid4().hex[:10]}" for _ in chunks]
    else:
        try:
            notebook_ids = await batch_upload(
                markdown_files=md_files,
                notebook_title=body.notebook_title,
                cookies_path=_COOKIES_PATH,
            )
        except NotebookUploadError as exc:
            raise HTTPException(status_code=exc.status_code, detail=str(exc)) from exc

    storage = _get_storage()
    passport_meta = storage.get_passport(body.passport_id) if body.passport_id else None
    for index, notebook_id in enumerate(notebook_ids, start=1):
        title = body.notebook_title if len(notebook_ids) == 1 else f"{body.notebook_title} — Part {index}"
        storage.register_notebook(
            notebook_id=notebook_id,
            title=title,
            output_dir=output_dir,
            source_paths=md_files,
            kind="master",
            passport_id=body.passport_id,
            taxonomy_version=passport_meta.get("taxonomy_version") if passport_meta else None,
            source_count=len(md_files),
        )

    return {"notebook_ids": notebook_ids, "source_count": len(md_files)}


@app.get("/passports/{passport_id}", tags=["convert"])
async def get_passport(passport_id: str) -> dict[str, Any]:
    storage = _get_storage()
    passport = storage.get_passport(passport_id)
    if passport is None:
        raise HTTPException(status_code=404, detail=f"Passport '{passport_id}' not found.")
    markdown = Path(passport["path"]).read_text(encoding="utf-8")
    return {
        **passport,
        "passport_markdown": markdown,
        "taxonomy": storage.get_passport_taxonomy(passport_id),
    }


@app.get("/passports/{passport_id}/taxonomy", tags=["convert"])
async def get_passport_taxonomy(passport_id: str) -> dict[str, Any]:
    taxonomy = _get_storage().get_passport_taxonomy(passport_id)
    if taxonomy is None:
        raise HTTPException(status_code=404, detail=f"Passport taxonomy for '{passport_id}' not found.")
    return taxonomy


@app.get("/notebooks/{notebook_id}", tags=["notebooklm"])
async def get_notebook_detail(notebook_id: str) -> dict[str, Any]:
    storage = _get_storage()
    notebook = storage.get_notebook(notebook_id)
    if notebook is None:
        raise HTTPException(status_code=404, detail=f"Notebook '{notebook_id}' not found.")
    passport_taxonomy = None
    if notebook.get("passport_id"):
        passport_taxonomy = storage.get_passport_taxonomy(notebook["passport_id"])
    artifacts = []
    for artifact in storage.list_notebook_artifacts(notebook_id):
        artifacts.append(
            {
                **artifact,
                "download_urls": {
                    fmt: f"/artifacts/{artifact['artifact_id']}/download?format={fmt}"
                    for fmt in artifact.get("files", {})
                },
            }
        )
    return {
        **notebook,
        "artifacts": artifacts,
        "taxonomy": _filter_taxonomy_for_notebook(passport_taxonomy, notebook),
    }


@app.get("/notebooks/{notebook_id}/sources", tags=["notebooklm"])
async def get_notebook_sources(notebook_id: str) -> dict[str, Any]:
    notebook = _get_storage().get_notebook(notebook_id)
    if notebook is None:
        raise HTTPException(status_code=404, detail=f"Notebook '{notebook_id}' not found.")
    return {
        "notebook_id": notebook_id,
        "sources": _extract_notebook_sources(notebook),
        "source_count": notebook.get("source_count", len(notebook.get("source_paths", []))),
    }


@app.post("/notebooks/derived", tags=["notebooklm"])
async def create_derived_notebook(body: DerivedNotebookRequest) -> dict[str, Any]:
    if not _DEMO_MODE and not is_authenticated(_COOKIES_PATH):
        raise HTTPException(status_code=401, detail="Not authenticated. Call /auth/login first.")
    if not body.selections:
        raise HTTPException(status_code=400, detail="Select at least one subcategory.")

    storage = _get_storage()
    passport = storage.get_passport(body.passport_id)
    if passport is None:
        raise HTTPException(status_code=404, detail=f"Passport '{body.passport_id}' not found.")
    taxonomy = storage.get_passport_taxonomy(body.passport_id)
    if taxonomy is None:
        raise HTTPException(status_code=404, detail="Passport taxonomy is unavailable.")

    parent_output_dir = _resolve_parent_output_dir(body.parent_output_dir, body.parent_notebook_id)
    conversations = _load_manifest_conversations(parent_output_dir)
    selected_ids = _resolve_selection_ids(taxonomy, body.selections)
    filtered_conversations = _filter_conversations_by_ids(conversations, selected_ids)
    if not filtered_conversations:
        raise HTTPException(status_code=404, detail="No conversations matched the selected subcategories.")

    derived_output_dir = _CONVERTED_DIR / f"derived_{uuid.uuid4().hex[:10]}"
    created_files = await asyncio.to_thread(
        _write_filtered_conversations,
        filtered_conversations,
        derived_output_dir,
    )
    await asyncio.to_thread(_save_conversion_manifest, derived_output_dir, filtered_conversations)

    if _DEMO_MODE:
        notebook_ids = [f"demo_{uuid.uuid4().hex[:10]}"]
    else:
        notebook_ids = await batch_upload(
            markdown_files=created_files,
            notebook_title=body.title,
            cookies_path=_COOKIES_PATH,
        )

    notebook_id = notebook_ids[0]
    notebook = storage.register_notebook(
        notebook_id=notebook_id,
        title=body.title,
        output_dir=derived_output_dir,
        source_paths=created_files,
        kind="thematic",
        passport_id=body.passport_id,
        taxonomy_version=passport.get("taxonomy_version"),
        parent_notebook_id=body.parent_notebook_id,
        selection_filters=[selection.model_dump() for selection in body.selections],
        source_count=len(created_files),
        conversation_ids=selected_ids,
    )
    return notebook


@app.post("/notebooks/{notebook_id}/artifacts", tags=["notebooklm"])
async def create_artifact_job(notebook_id: str, body: ArtifactJobRequest) -> dict[str, Any]:
    if not _DEMO_MODE and not is_authenticated(_COOKIES_PATH):
        raise HTTPException(status_code=401, detail="Not authenticated.")

    notebook = _get_storage().get_notebook(notebook_id)
    if notebook is None:
        raise HTTPException(status_code=404, detail=f"Notebook '{notebook_id}' not found.")

    requested_types = _normalize_artifact_types(body.types)
    record = jobs.create(
        notebook_id=notebook_id,
        artifact_types=requested_types,
        runner=lambda job_id: _run_artifact_job(job_id, notebook_id, requested_types),
    )
    return record


@app.get("/jobs/{job_id}", tags=["notebooklm"])
async def get_job(job_id: str) -> dict[str, Any]:
    job = jobs.get(job_id)
    if job is None:
        raise HTTPException(status_code=404, detail=f"Job '{job_id}' not found.")
    return job


@app.get("/artifacts/{artifact_id}", tags=["notebooklm"])
async def get_artifact(artifact_id: str) -> dict[str, Any]:
    artifact = _get_storage().get_artifact(artifact_id)
    if artifact is None:
        raise HTTPException(status_code=404, detail=f"Artifact '{artifact_id}' not found.")

    download_urls = {
        fmt: f"/artifacts/{artifact_id}/download?format={fmt}"
        for fmt in artifact.get("files", {})
    }
    return {
        **artifact,
        "download_urls": download_urls,
    }


@app.get("/artifacts/{artifact_id}/download", tags=["notebooklm"])
async def download_artifact(artifact_id: str, format: str) -> FileResponse:
    path = _get_storage().get_artifact_file(artifact_id, format)
    if path is None:
        raise HTTPException(
            status_code=404,
            detail=f"Artifact '{artifact_id}' does not have format '{format}'.",
        )
    return FileResponse(path=str(path), filename=path.name)


@app.post("/artifacts/{artifact_id}/export", tags=["notebooklm"])
async def export_artifact(artifact_id: str, body: ExportArtifactRequest) -> FileResponse:
    return await download_artifact(artifact_id, body.format)


@app.patch("/notebooks/{notebook_id}/slides/{slide_index}", tags=["notebooklm"])
async def revise_notebook_slide(
    notebook_id: str,
    slide_index: int,
    body: ReviseSlideRequest,
) -> dict[str, Any]:
    if not _DEMO_MODE and not is_authenticated(_COOKIES_PATH):
        raise HTTPException(status_code=401, detail="Not authenticated.")

    if _DEMO_MODE:
        artifact = _get_storage().get_artifact(body.artifact_id)
        if artifact is None:
            raise HTTPException(status_code=404, detail=f"Artifact '{body.artifact_id}' not found.")
        slides = artifact.get("preview", {}).get("slides", [])
        if slide_index < 0 or slide_index >= len(slides):
            raise HTTPException(status_code=400, detail="Slide index is out of range.")
        updated = {
            **slides[slide_index],
            "content": f"{slides[slide_index].get('content', '')}\n\nRevision request: {body.revision_prompt}",
        }
    else:
        updated = await revise_slide(
            notebook_id=notebook_id,
            artifact_id=body.artifact_id,
            slide_index=slide_index,
            revision_prompt=body.revision_prompt,
            cookies_path=_COOKIES_PATH,
        )

    artifact = _get_storage().get_artifact(body.artifact_id)
    if artifact is None:
        raise HTTPException(status_code=404, detail=f"Artifact '{body.artifact_id}' not found.")

    slides = artifact.get("preview", {}).get("slides", [])
    if slide_index < 0 or slide_index >= len(slides):
        raise HTTPException(status_code=400, detail="Slide index is out of range.")
    slides[slide_index] = updated

    storage = _get_storage()
    storage.add_artifact_file(
        body.artifact_id,
        "json",
        content=json.dumps({"artifact_id": body.artifact_id, "slides": slides}, indent=2),
        filename="slides.json",
        preview={"artifact_id": body.artifact_id, "slides": slides},
    )
    storage.add_artifact_file(
        body.artifact_id,
        "md",
        content=_render_slides_markdown(slides),
        filename="slides.md",
    )
    storage.add_artifact_file(
        body.artifact_id,
        "csv",
        content=build_csv(derive_slide_table_rows(slides)),
        filename="tables.csv",
    )
    storage.add_artifact_file(
        body.artifact_id,
        "table_md",
        content=build_markdown_table(derive_slide_table_rows(slides)),
        filename="tables.md",
    )
    _maybe_add_pptx_export(storage, body.artifact_id, slides)

    return {"notebook_id": notebook_id, "slide_index": slide_index, "slide": updated}


@app.post("/search/index", tags=["search"])
async def build_search_index(body: SearchIndexRequest) -> dict[str, Any]:
    from byegpt.indexer import VectorIndexer

    input_dir = Path(body.input_dir)
    if not input_dir.exists():
        raise HTTPException(status_code=404, detail=f"Input dir '{input_dir}' not found.")

    indexer = VectorIndexer(_SEARCH_DIR)
    indexed_files = await asyncio.to_thread(indexer.index_directory, input_dir)
    return {
        "status": "indexed",
        "input_dir": str(input_dir),
        "db_path": str(_SEARCH_DIR),
        "indexed_files": indexed_files,
        "document_count": indexer.count(),
    }


@app.get("/search/index/status", tags=["search"])
async def search_index_status() -> dict[str, Any]:
    from byegpt.indexer import VectorIndexer

    indexer = VectorIndexer(_SEARCH_DIR)
    return {
        "db_path": str(_SEARCH_DIR),
        "document_count": indexer.count(),
        "ready": indexer.count() > 0,
    }


@app.post("/search/query", tags=["search"])
async def search_query(body: SearchQueryRequest) -> dict[str, Any]:
    from byegpt.indexer import VectorIndexer

    indexer = VectorIndexer(_SEARCH_DIR)
    return {
        "results": indexer.query(body.text, n_results=body.n_results),
    }


@app.get("/notebooks/{notebook_id}/export", tags=["notebooklm"])
async def export_notebook_bundle(notebook_id: str) -> FileResponse:
    try:
        archive_path = _get_storage().create_notebook_export(notebook_id)
    except FileNotFoundError as exc:
        raise HTTPException(status_code=404, detail=str(exc)) from exc
    return FileResponse(path=str(archive_path), filename=archive_path.name)


async def _convert_archive(
    input_path: Path,
    output_dir: Path,
    max_size_mb: float,
    include_thinking: bool,
    include_attachments: bool,
) -> dict[str, Any]:
    try:
        result = await asyncio.to_thread(
            convert_with_anchors,
            input_path=input_path,
            output_dir=output_dir,
            max_size_mb=max_size_mb,
            include_thinking=include_thinking,
            include_attachments=include_attachments,
        )
    except Exception as exc:
        logger.exception("Conversion failed for %s", input_path)
        raise HTTPException(status_code=500, detail=f"Conversion failed: {exc}") from exc

    from byegpt.parser import load_conversations

    conversations, zf = await asyncio.to_thread(load_conversations, input_path)
    if zf is not None:
        zf.close()

    await asyncio.to_thread(_save_conversion_manifest, output_dir, conversations)
    topic_laboratory = build_topic_laboratory(conversations)
    taxonomy = build_taxonomy(conversations)

    return {
        "output_dir": str(output_dir),
        "files_created": len(result.created_files),
        "attachment_count": result.attachment_count,
        "conversation_count": result.conversation_count,
        "file_paths": [str(path) for path in result.created_files],
        "topic_laboratory": topic_laboratory,
        "conversation_map": [
            {
                "conversation_id": conversation_uid(conversation, index),
                "title": conversation.get("title") or "Untitled conversation",
            }
            for index, conversation in enumerate(conversations)
        ],
        "taxonomy": taxonomy,
    }


async def _run_artifact_job(
    job_id: str,
    notebook_id: str,
    artifact_types: list[str],
) -> dict[str, Any]:
    storage = _get_storage()
    tasks = [asyncio.create_task(_generate_and_store_artifact(notebook_id, artifact_type)) for artifact_type in artifact_types]
    artifacts = await asyncio.gather(*tasks)
    artifact_ids = [artifact["artifact_id"] for artifact in artifacts]
    return {
        "job_id": job_id,
        "artifact_ids": artifact_ids,
        "artifacts": [
            {
                **storage.get_artifact(artifact_id),
                "download_urls": {
                    fmt: f"/artifacts/{artifact_id}/download?format={fmt}"
                    for fmt in storage.get_artifact(artifact_id).get("files", {})
                },
            }
            for artifact_id in artifact_ids
            if storage.get_artifact(artifact_id) is not None
        ],
    }


async def _generate_and_store_artifact(notebook_id: str, artifact_type: str) -> dict[str, Any]:
    storage = _get_storage()
    if _DEMO_MODE:
        return await _generate_demo_artifact(notebook_id, artifact_type)

    if artifact_type == "mind_map":
        result = await generate_mind_map(notebook_id=notebook_id, cookies_path=_COOKIES_PATH)
        metadata = storage.create_artifact(
            notebook_id=notebook_id,
            artifact_type="mind_map",
            upstream_artifact_id=result["artifact_id"],
        )
        storage.add_artifact_file(
            metadata["artifact_id"],
            "json",
            content=json.dumps(result["data"], indent=2),
            filename="mind_map.json",
            preview=result["data"],
        )
        return metadata

    if artifact_type == "audio":
        result = await generate_audio_overview(notebook_id=notebook_id, cookies_path=_COOKIES_PATH)
        metadata = storage.create_artifact(
            notebook_id=notebook_id,
            artifact_type="audio",
            upstream_artifact_id=result["artifact_id"],
        )
        storage.add_artifact_file(
            metadata["artifact_id"],
            "mp4",
            content=result["bytes"],
            filename="audio.mp4",
        )
        return metadata

    if artifact_type == "slides":
        result = await generate_slides(notebook_id=notebook_id, cookies_path=_COOKIES_PATH)
        metadata = storage.create_artifact(
            notebook_id=notebook_id,
            artifact_type="slides",
            upstream_artifact_id=result["artifact_id"],
        )
        slides = result["slides"]
        storage.add_artifact_file(
            metadata["artifact_id"],
            "json",
            content=json.dumps({"artifact_id": result["artifact_id"], "slides": slides}, indent=2),
            filename="slides.json",
            preview={"artifact_id": result["artifact_id"], "slides": slides},
        )
        storage.add_artifact_file(
            metadata["artifact_id"],
            "md",
            content=_render_slides_markdown(slides),
            filename="slides.md",
        )
        rows = derive_slide_table_rows(slides)
        storage.add_artifact_file(
            metadata["artifact_id"],
            "csv",
            content=build_csv(rows),
            filename="tables.csv",
        )
        storage.add_artifact_file(
            metadata["artifact_id"],
            "table_md",
            content=build_markdown_table(rows),
            filename="tables.md",
        )
        if result.get("pptx_bytes"):
            storage.add_artifact_file(
                metadata["artifact_id"],
                "pptx",
                content=result["pptx_bytes"],
                filename="slides.pptx",
            )
        if result.get("pdf_bytes"):
            storage.add_artifact_file(
                metadata["artifact_id"],
                "pdf",
                content=result["pdf_bytes"],
                filename="slides.pdf",
            )
        return metadata

    if artifact_type == "quiz":
        result = await generate_quiz(notebook_id=notebook_id, cookies_path=_COOKIES_PATH)
        metadata = storage.create_artifact(
            notebook_id=notebook_id,
            artifact_type="quiz",
            upstream_artifact_id=result["artifact_id"],
        )
        quiz = result["quiz"]
        storage.add_artifact_file(
            metadata["artifact_id"],
            "json",
            content=json.dumps(quiz, indent=2),
            filename="quiz.json",
            preview=quiz,
        )
        storage.add_artifact_file(
            metadata["artifact_id"],
            "md",
            content=_render_quiz_markdown(quiz),
            filename="quiz.md",
        )
        return metadata

    raise HTTPException(status_code=400, detail=f"Unsupported artifact type '{artifact_type}'.")


async def _generate_demo_artifact(notebook_id: str, artifact_type: str) -> dict[str, Any]:
    notebook = _get_storage().get_notebook(notebook_id)
    if notebook is None:
        raise HTTPException(status_code=404, detail=f"Notebook '{notebook_id}' not found.")

    titles = _extract_notebook_titles(notebook)
    storage = _get_storage()

    if artifact_type == "mind_map":
        nodes = [{"id": "root", "label": notebook.get("title", "Notebook"), "group": "root"}]
        links: list[dict[str, str]] = []
        for index, title in enumerate(titles[:12], start=1):
            node_id = f"topic_{index}"
            nodes.append({"id": node_id, "label": title, "group": "conversation"})
            links.append({"source": "root", "target": node_id})
        metadata = storage.create_artifact(notebook_id=notebook_id, artifact_type="mind_map")
        payload = {"nodes": nodes, "links": links}
        storage.add_artifact_file(
            metadata["artifact_id"],
            "json",
            content=json.dumps(payload, indent=2),
            filename="mind_map.json",
            preview=payload,
        )
        return metadata

    if artifact_type == "audio":
        metadata = storage.create_artifact(notebook_id=notebook_id, artifact_type="audio")
        storage.add_artifact_file(
            metadata["artifact_id"],
            "wav",
            content=_build_demo_wav(),
            filename="audio.wav",
        )
        return metadata

    if artifact_type == "slides":
        slides = [
            {
                "title": f"Insight {index}",
                "content": f"Key theme from your archive: {title}",
            }
            for index, title in enumerate(titles[:5], start=1)
        ] or [{"title": "Insight 1", "content": "Upload a richer export to generate more detailed slides."}]
        metadata = storage.create_artifact(notebook_id=notebook_id, artifact_type="slides")
        storage.add_artifact_file(
            metadata["artifact_id"],
            "json",
            content=json.dumps({"artifact_id": metadata["artifact_id"], "slides": slides}, indent=2),
            filename="slides.json",
            preview={"artifact_id": metadata["artifact_id"], "slides": slides},
        )
        storage.add_artifact_file(
            metadata["artifact_id"],
            "md",
            content=_render_slides_markdown(slides),
            filename="slides.md",
        )
        rows = derive_slide_table_rows(slides)
        storage.add_artifact_file(
            metadata["artifact_id"],
            "csv",
            content=build_csv(rows),
            filename="tables.csv",
        )
        storage.add_artifact_file(
            metadata["artifact_id"],
            "table_md",
            content=build_markdown_table(rows),
            filename="tables.md",
        )
        _maybe_add_pptx_export(storage, metadata["artifact_id"], slides)
        return metadata

    if artifact_type == "quiz":
        questions = [
            {
                "question": f"What topic did the archive cover around '{title}'?",
                "answer": title,
            }
            for title in titles[:5]
        ] or [{"question": "What does demo mode do?", "answer": "It lets you test the product without NotebookLM auth."}]
        quiz = {"title": "Archive Quiz", "questions": questions}
        metadata = storage.create_artifact(notebook_id=notebook_id, artifact_type="quiz")
        storage.add_artifact_file(
            metadata["artifact_id"],
            "json",
            content=json.dumps(quiz, indent=2),
            filename="quiz.json",
            preview=quiz,
        )
        storage.add_artifact_file(
            metadata["artifact_id"],
            "md",
            content=_render_quiz_markdown(quiz),
            filename="quiz.md",
        )
        return metadata

    raise HTTPException(status_code=400, detail=f"Unsupported artifact type '{artifact_type}'.")


def _normalize_artifact_types(types: list[str]) -> list[str]:
    allowed = {"mind_map", "audio", "slides", "quiz"}
    normalized: list[str] = []
    for artifact_type in types:
        if artifact_type not in allowed:
            raise HTTPException(status_code=400, detail=f"Unsupported artifact type '{artifact_type}'.")
        if artifact_type not in normalized:
            normalized.append(artifact_type)
    return normalized or ["mind_map", "audio", "slides", "quiz"]


def _render_slides_markdown(slides: list[dict[str, Any]]) -> str:
    lines = ["# Slides", ""]
    for index, slide in enumerate(slides, start=1):
        lines.append(f"## {index}. {slide.get('title', f'Slide {index}')}")
        lines.append("")
        lines.append(str(slide.get("content", "")).strip())
        lines.append("")
    return "\n".join(lines).strip() + "\n"


def _render_quiz_markdown(quiz: dict[str, Any]) -> str:
    lines = [f"# {quiz.get('title', 'Quiz')}", ""]
    for index, question in enumerate(quiz.get("questions", []), start=1):
        lines.append(f"## Question {index}")
        lines.append("")
        lines.append(str(question.get("question", "")))
        lines.append("")
        answer = question.get("answer")
        if answer:
            lines.append(f"Answer: {answer}")
            lines.append("")
    return "\n".join(lines).strip() + "\n"


def _maybe_add_pptx_export(storage: StorageManager, artifact_id: str, slides: list[dict[str, Any]]) -> None:
    try:
        from pptx import Presentation  # type: ignore[import]
    except ImportError:
        return

    presentation = Presentation()
    layout = presentation.slide_layouts[1]
    for slide in slides:
        ppt_slide = presentation.slides.add_slide(layout)
        ppt_slide.shapes.title.text = slide.get("title", "Slide")
        ppt_slide.placeholders[1].text = slide.get("content", "")

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        tmp_path = Path(tmp.name)
    try:
        presentation.save(str(tmp_path))
        storage.add_artifact_file(
            artifact_id,
            "pptx",
            content=tmp_path.read_bytes(),
            filename="slides.pptx",
        )
    finally:
        tmp_path.unlink(missing_ok=True)


def _extract_notebook_titles(notebook: dict[str, Any]) -> list[str]:
    titles: list[str] = []
    for source_path in notebook.get("source_paths", []):
        path = Path(source_path)
        if not path.exists():
            continue
        content = path.read_text(encoding="utf-8")
        for line in content.splitlines():
            if line.startswith('title: "'):
                titles.append(line.replace('title: "', "").rstrip('"'))
                break
    return titles


def _extract_notebook_sources(notebook: dict[str, Any]) -> list[dict[str, Any]]:
    sources = []
    for source_path in notebook.get("source_paths", []):
        path = Path(source_path)
        if not path.exists():
            continue
        title = path.stem
        conversation_id = None
        for line in path.read_text(encoding="utf-8", errors="replace").splitlines():
            if line.startswith("<!-- source: https://chatgpt.com/c/"):
                conversation_id = line.rsplit("/", 1)[-1].replace(" -->", "").strip()
                continue
            if line.startswith('title: "'):
                title = line.replace('title: "', "").rstrip('"')
                break
        sources.append(
            {
                "path": str(path),
                "title": title,
                "conversation_id": conversation_id,
            }
        )
    return sources


def _conversion_manifest_path(output_dir: Path) -> Path:
    return Path(output_dir) / ".byegpt_manifest.json"


def _save_conversion_manifest(output_dir: Path, conversations: list[dict[str, Any]]) -> None:
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)
    manifest = {
        "conversation_count": len(conversations),
        "conversations": conversations,
    }
    _conversion_manifest_path(output_dir).write_text(
        json.dumps(manifest),
        encoding="utf-8",
    )


def _load_manifest_conversations(output_dir: Path) -> list[dict[str, Any]]:
    manifest_path = _conversion_manifest_path(output_dir)
    if not manifest_path.exists():
        raise HTTPException(
            status_code=404,
            detail=f"Conversion manifest not found for '{output_dir}'. Re-run conversion first.",
        )
    payload = json.loads(manifest_path.read_text(encoding="utf-8"))
    return payload.get("conversations", [])


def _resolve_parent_output_dir(
    parent_output_dir: str | None,
    parent_notebook_id: str | None,
) -> Path:
    if parent_output_dir:
        return Path(parent_output_dir)
    if parent_notebook_id:
        notebook = _get_storage().get_notebook(parent_notebook_id)
        if notebook is None:
            raise HTTPException(status_code=404, detail=f"Parent notebook '{parent_notebook_id}' not found.")
        return Path(notebook["output_dir"])
    raise HTTPException(status_code=400, detail="Provide parent_output_dir or parent_notebook_id.")


def _resolve_selection_ids(
    taxonomy: dict[str, Any],
    selections: list[NotebookSelection],
) -> list[str]:
    matched: list[str] = []
    for selection in selections:
        category_slug = _slugify_token(selection.category)
        subcategory_slug = _slugify_token(selection.subcategory)
        for category in taxonomy.get("categories", []):
            if category.get("slug") != category_slug:
                continue
            for subcategory in category.get("subcategories", []):
                if subcategory.get("slug") == subcategory_slug:
                    matched.extend(subcategory.get("conversation_ids", []))
    # Deduplicate, preserve order.
    return list(dict.fromkeys(matched))


def _filter_conversations_by_ids(
    conversations: list[dict[str, Any]],
    selected_ids: list[str],
) -> list[dict[str, Any]]:
    selected = set(selected_ids)
    filtered = []
    for index, conversation in enumerate(conversations):
        if conversation_uid(conversation, index) in selected:
            filtered.append(conversation)
    return filtered


def _slugify_token(value: str) -> str:
    return re.sub(r"[^a-z0-9]+", "_", value.strip().lower()).strip("_")


def _write_filtered_conversations(
    conversations: list[dict[str, Any]],
    output_dir: Path,
) -> list[Path]:
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)
    created_files: list[Path] = []
    for index, conversation in enumerate(conversations, start=1):
        title = (conversation.get("title") or f"Conversation {index}").strip()
        safe_title = re.sub(r"[^\w.-]+", "_", title).strip("_").lower() or f"conversation_{index}"
        path = output_dir / f"{index:03d}_{safe_title[:80]}.md"
        markdown = format_conversation(conversation, attachment_map={}, include_thinking=True)
        path.write_text(inject_context_anchors(markdown, conversation), encoding="utf-8")
        created_files.append(path)
    return created_files


def _filter_taxonomy_for_notebook(
    taxonomy: dict[str, Any] | None,
    notebook: dict[str, Any],
) -> dict[str, Any] | None:
    if taxonomy is None:
        return None
    selected_ids = notebook.get("conversation_ids") or []
    if not selected_ids:
        return taxonomy
    selected = set(selected_ids)
    filtered_categories = []
    for category in taxonomy.get("categories", []):
        filtered_subcategories = []
        category_count = 0
        for subcategory in category.get("subcategories", []):
            kept_ids = [conversation_id for conversation_id in subcategory.get("conversation_ids", []) if conversation_id in selected]
            if not kept_ids:
                continue
            filtered_subcategory = {
                **subcategory,
                "conversation_ids": kept_ids,
                "conversations": [
                    conversation
                    for conversation in subcategory.get("conversations", [])
                    if conversation.get("conversation_id") in selected
                ],
                "count": len(kept_ids),
            }
            filtered_subcategories.append(filtered_subcategory)
            category_count += len(kept_ids)
        if filtered_subcategories:
            filtered_categories.append(
                {
                    **category,
                    "count": category_count,
                    "subcategories": filtered_subcategories,
                }
            )
    return {
        **taxonomy,
        "categories": filtered_categories,
        "total_conversations": len(selected_ids),
    }


def _build_demo_wav() -> bytes:
    buffer = BytesIO()
    with wave.open(buffer, "wb") as wav_file:
        wav_file.setnchannels(1)
        wav_file.setsampwidth(2)
        wav_file.setframerate(22050)
        frames = bytearray()
        for index in range(22050):
            amplitude = int(12000 * __import__("math").sin(2 * __import__("math").pi * 440 * index / 22050))
            frames.extend(int(amplitude).to_bytes(2, byteorder="little", signed=True))
        wav_file.writeframes(bytes(frames))
    return buffer.getvalue()


def _load_tus_upload_info(upload_id: str) -> TusUploadInfo:
    info_path = _TUSD_UPLOAD_DIR / f"{upload_id}.info"
    if not info_path.exists():
        raise HTTPException(status_code=404, detail=f"Tus upload '{upload_id}' not found.")

    try:
        payload = json.loads(info_path.read_text(encoding="utf-8"))
        return TusUploadInfo.model_validate(payload)
    except (json.JSONDecodeError, ValueError) as exc:
        raise HTTPException(status_code=400, detail="Tus upload metadata is invalid.") from exc


def _normalize_tusd_metadata(value: str) -> str:
    if not value:
        return value

    try:
        decoded = base64.b64decode(value, validate=True).decode("utf-8")
    except (binascii.Error, UnicodeDecodeError):
        return value

    if decoded and all(char.isprintable() or char.isspace() for char in decoded):
        return decoded
    return value
